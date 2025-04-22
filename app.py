from flask import Flask, render_template, request, redirect, url_for, flash, session
import os
from werkzeug.utils import secure_filename
import tempfile
from pptx import Presentation
import google.generativeai as genai
import secrets

app = Flask(__name__)
app.secret_key = secrets.token_hex(16)

# Configuration
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'ppt', 'pptx'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max upload

# Create upload folder if it doesn't exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Set your Google Gemini API key
# For production, use environment variables instead
GEMINI_API_KEY = "your api key"  # Replace with your actual Gemini API key
genai.configure(api_key=GEMINI_API_KEY)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_ppt(file_path):
    """Extract text content from a PowerPoint file"""
    ppt = Presentation(file_path)
    text_content = []
    
    for slide in ppt.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
        text_content.append(" ".join(slide_text))
    
    return "\n\n".join(text_content)

def generate_summary(text):
    """Generate a summary using Google Gemini API"""
    if not text.strip():
        return "No text content found in the presentation."
    
    try:
        # Use a specific model that is available according to the error message
        model = genai.GenerativeModel("models/gemini-1.5-flash")
        
        # Create a prompt for summarization
        prompt = f"""
        You are a helpful assistant that summarizes PowerPoint presentations.
        Please summarize the following PowerPoint content in a concise way, highlighting the main points and key takeaways:
        
        {text}
        """
        
        # Generate the response
        response = model.generate_content(prompt)
        
        return response.text
    except Exception as e:
        # Try fallback to another model if the first one fails
        try:
            fallback_model = genai.GenerativeModel("models/gemini-1.5-pro")
            response = fallback_model.generate_content(prompt)
            return response.text
        except Exception as fallback_e:
            return f"Error generating summary with primary model: {str(e)}\n\nError with fallback model: {str(fallback_e)}"

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        flash('No file part')
        return redirect(request.url)
    
    file = request.files['file']
    
    if file.filename == '':
        flash('No selected file')
        return redirect(request.url)
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        # Extract text content
        text_content = extract_text_from_ppt(file_path)
        
        # Store in session
        session['file_path'] = file_path
        session['text_content'] = text_content
        session['filename'] = filename
        
        return redirect(url_for('preview'))
    else:
        flash('Invalid file type. Please upload a PPT or PPTX file.')
        return redirect(request.url)

@app.route('/preview')
def preview():
    if 'filename' not in session:
        return redirect(url_for('index'))
    
    return render_template('preview.html', 
                          filename=session.get('filename'), 
                          text_preview=session.get('text_content')[:500] + '...' if len(session.get('text_content', '')) > 500 else session.get('text_content'))

@app.route('/summarize', methods=['POST'])
def summarize():
    if 'text_content' not in session:
        flash('No presentation content found')
        return redirect(url_for('index'))
    
    text_content = session.get('text_content')
    summary = generate_summary(text_content)
    
    return render_template('summary.html', 
                          filename=session.get('filename'),
                          summary=summary)

if __name__ == '__main__':
    app.run(debug=True)
